import streamlit as st
from transformers import pipeline
from sentence_transformers import SentenceTransformer
from bs4 import BeautifulSoup
import requests
import faiss
import numpy as np
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract

# Initialize models
@st.cache_resource
def load_models():
    summarizer = pipeline("summarization", model="facebook/bart-large-cnn")
    embedder = SentenceTransformer("sentence-transformers/multi-qa-MiniLM-L6-cos-v1")
    llm = pipeline("text2text-generation", model="google/flan-t5-large")
    return summarizer, embedder, llm

summarizer, embedder, llm = load_models()

# FAISS vector store
class VectorStore:
    def __init__(self):
        self.index = None
        self.sentences = []

    def add_sentences(self, sentences):
        if not sentences:
            return
        embeddings = embedder.encode(sentences, convert_to_tensor=False)
        embeddings = np.array(embeddings, dtype="float32")
        if self.index is None:
            self.index = faiss.IndexFlatL2(embeddings.shape[1])
        self.index.add(embeddings)
        self.sentences.extend(sentences)

    def search(self, query, top_k=3):
        if not self.index or self.index.ntotal == 0:
            return []
        query_embedding = np.array([embedder.encode(query)], dtype="float32")
        distances, indices = self.index.search(query_embedding, top_k)
        results = [
            (self.sentences[i], distances[0][j])
            for j, i in enumerate(indices[0])
            if 0 <= i < len(self.sentences)
        ]
        return results

vector_store = VectorStore()

# Extract text from the given URL
def extract_text_from_url(url):
    try:
        response = requests.get(url)
        response.raise_for_status()
        soup = BeautifulSoup(response.text, "html.parser")
        paragraphs = soup.find_all("p")
        text = " ".join([para.get_text() for para in paragraphs])
        return text.strip()
    except Exception as e:
        return f"Error fetching data from URL: {e}"

# Extract text from a document
def extract_text_from_document(file):
    text = ""
    try:
        if file.type == "application/pdf":
            pdf_reader = PdfReader(file)
            text = " ".join(page.extract_text() for page in pdf_reader.pages if page.extract_text())
        elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            doc = Document(file)
            text = " ".join(para.text for para in doc.paragraphs if para.text)
        elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            ppt = Presentation(file)
            text = " ".join(
                shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text") and shape.text
            )
        elif file.type == "text/plain":
            text = file.read().decode("utf-8")
        elif file.type.startswith("image/"):
            image = Image.open(file)
            text = pytesseract.image_to_string(image)
        else:
            raise ValueError("Unsupported file type.")
        return text.strip()
    except Exception as e:
        return f"Error reading document: {e}"

# Split text into sentences
def split_text_into_sentences(text):
    sentences = text.replace("\n", " ").split(". ")
    return [sentence.strip() for sentence in sentences if sentence.strip()]

# Modify query
def modify_query(query, relevant_texts):
    context = "\n".join(relevant_texts)
    modified_query = f"Context:\n{context}\n\nQuestion: {query}"
    return modified_query

# Get LLM answer
def get_llm_answer(query, relevant_texts):
    modified_query = modify_query(query, relevant_texts)
    response = llm(modified_query, max_length=200, do_sample=False)
    return response[0]['generated_text']

# Streamlit App
st.title("SmartChat")
st.markdown("### Text Summarizer and Enhanced QA")
st.markdown("### Select an option: Chat with Links or Chat with Docs")

# Select Chat Mode
chat_mode = st.radio("Select Mode:", ["Chat with Links", "Chat with Docs"])

if chat_mode == "Chat with Links":
    url = st.text_input("Enter the URL:")

    if url:
        with st.spinner("Extracting text from the URL..."):
            extracted_text = extract_text_from_url(url)
            if extracted_text.startswith("Error"):
                st.error(extracted_text)
            else:
                st.success("Text extracted successfully!")
                st.text_area("Extracted Text", extracted_text, height=200)

                sentences = split_text_into_sentences(extracted_text)
                vector_store.add_sentences(sentences)
                st.success("Sentences added to the vector store!")

                if st.button("Summarize Text"):
                    with st.spinner("Generating summary..."):
                        try:
                            summary = []
                            for chunk in sentences[:10]:
                                chunk_summary = summarizer(chunk, max_length=150, min_length=40, do_sample=False)
                                summary.append(chunk_summary[0]["summary_text"])

                            final_summary = " ".join(summary)
                            st.markdown("### Summary:")
                            st.write(final_summary)
                        except Exception as e:
                            st.error(f"Error during summarization: {e}")

                query = st.text_input("Ask a question or enter a search query:")
                if query:
                    with st.spinner("Searching the vector store..."):
                        results = vector_store.search(query, top_k=3)
                        relevant_texts = [result[0] for result in results]
                        if relevant_texts:
                            with st.spinner("Generating answer..."):
                                try:
                                    answer = get_llm_answer(query, relevant_texts)
                                    st.markdown("### Answer:")
                                    st.write(answer)
                                except Exception as e:
                                    st.error(f"Error during answer generation: {e}")
                        else:
                            st.warning("No relevant results found. Try rephrasing your query.")

elif chat_mode == "Chat with Docs":
    uploaded_file = st.file_uploader("Upload a document (PDF, Word, PPT, TXT, or Image):")

    if uploaded_file:
        with st.spinner("Extracting text from the document..."):
            extracted_text = extract_text_from_document(uploaded_file)
            if extracted_text.startswith("Error"):
                st.error(extracted_text)
            elif not extracted_text:
                st.warning("The uploaded file contains no extractable text.")
            else:
                st.success("Text extracted successfully!")
                st.text_area("Extracted Text", extracted_text, height=200)

                sentences = split_text_into_sentences(extracted_text)
                if sentences:
                    vector_store.add_sentences(sentences)
                    st.success("Sentences added to the vector store!")
                else:
                    st.warning("No valid sentences were found in the extracted text.")

                if st.button("Summarize Text"):
                    with st.spinner("Generating summary..."):
                        try:
                            summary = []
                            for chunk in sentences[:10]:
                                chunk_summary = summarizer(chunk, max_length=150, min_length=40, do_sample=False)
                                summary.append(chunk_summary[0]["summary_text"])

                            final_summary = " ".join(summary)
                            st.markdown("### Summary:")
                            st.write(final_summary)
                        except Exception as e:
                            st.error(f"Error during summarization: {e}")

                query = st.text_input("Ask a question or enter a search query:")
                if query:
                    with st.spinner("Searching the vector store..."):
                        results = vector_store.search(query, top_k=3)
                        relevant_texts = [result[0] for result in results]
                        if relevant_texts:
                            with st.spinner("Generating answer..."):
                                try:
                                    answer = get_llm_answer(query, relevant_texts)
                                    st.markdown("### Answer:")
                                    st.write(answer)
                                except Exception as e:
                                    st.error(f"Error during answer generation: {e}")
                        else:
                            st.warning("No relevant results found. Try rephrasing your query.")
