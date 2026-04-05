# app/services/document_processor.py
import fitz  # PyMuPDF
import docx
from PIL import Image
import pytesseract
from typing import List, Dict, Any
import hashlib
import uuid
import io
from app.models.schemas import DocumentChunk, DocumentMetadata, DocumentType
import markdown
from pptx import Presentation
from bs4 import BeautifulSoup
import html2text

class DocumentProcessor:
    def __init__(self):
        self.supported_formats = {
            '.pdf': self._process_pdf,
            '.docx': self._process_docx,
            '.txt': self._process_txt,
            '.png': self._process_image,
            '.jpg': self._process_image,
            '.jpeg': self._process_image,
            '.md': self._process_markdown,      # New
            '.markdown': self._process_markdown, # New
            '.pptx': self._process_powerpoint,  # New
            '.ppt': self._process_powerpoint,   # New (legacy)
            '.html': self._process_html,        # New
            '.htm': self._process_html          # New
        }
    
    # for markdown
    async def _process_markdown(self, file_path: str) -> str:
        """Extract text from Markdown files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                md_content = file.read()
            
            # Convert markdown to HTML, then extract plain text
            html_content = markdown.markdown(md_content, extensions=['tables', 'fenced_code'])
            
            # Convert HTML to plain text
            h = html2text.HTML2Text()
            h.ignore_links = False
            h.ignore_images = False
            plain_text = h.handle(html_content)
            
            return plain_text
            
        except Exception as e:
            raise Exception(f"Error processing Markdown file: {str(e)}")

    # for ppt
    async def _process_powerpoint(self, file_path: str) -> str:
        """Extract text from PowerPoint files"""
        try:
            presentation = Presentation(file_path)
            text = ""
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                text += f"\n--- Slide {slide_num} ---\n"
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text += shape.text + "\n"
                    
                    # Extract text from tables
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                row_text.append(cell.text.strip())
                            text += " | ".join(row_text) + "\n"
                    
                    # Extract text from text frames
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text += paragraph.text + "\n"
                
                # Extract notes
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    if notes_slide.notes_text_frame:
                        text += f"\nNotes: {notes_slide.notes_text_frame.text}\n"
            
            return text
            
        except Exception as e:
            raise Exception(f"Error processing PowerPoint file: {str(e)}")

    # for html
    async def _process_html(self, file_path: str) -> str:
        """Extract text from HTML files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                html_content = file.read()
            
            # Method 1: Using html2text for better formatting
            h = html2text.HTML2Text()
            h.ignore_links = False
            h.ignore_images = False
            h.body_width = 0  # Don't wrap lines
            text = h.handle(html_content)
            
            # Method 2: Fallback using BeautifulSoup for cleaner extraction
            if not text.strip():
                soup = BeautifulSoup(html_content, 'html.parser')
                
                # Remove script and style elements
                for script in soup(["script", "style", "nav", "header", "footer"]):
                    script.decompose()
                
                # Extract text
                text = soup.get_text(separator='\n', strip=True)
            
            return text
            
        except Exception as e:
            raise Exception(f"Error processing HTML file: {str(e)}")

    def _validate_file_format(self, file_path: str, filename: str) -> bool:
        """Validate file format and content"""
        file_extension = filename.lower().split('.')[-1]
        
        try:
            if file_extension in ['pptx', 'ppt']:
                # Validate PowerPoint file
                Presentation(file_path)
            elif file_extension in ['html', 'htm']:
                # Validate HTML file
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                    if not content.strip():
                        return False
            elif file_extension in ['md', 'markdown']:
                # Validate Markdown file
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
                    if not content.strip():
                        return False
            
            return True
            
        except Exception as e:
            return False
            # raise ValueError(f"File validation failed for {filename}: {str(e)}")

    # for docx
    async def process_document(self, file_path: str, filename: str) -> List[DocumentChunk]:
        """Process document and return chunks"""
        file_extension = filename.lower().split('.')[-1]
        file_extension = f".{file_extension}"
        
        if file_extension not in self.supported_formats:
            raise ValueError(f"Unsupported file format: {file_extension}")
        
        if not self._validate_file_format(file_path, filename):
            raise ValueError(f"Invalid or corrupted file: {filename}")


        # Extract text content
        content = await self.supported_formats[file_extension](file_path)
        
        # Create chunks
        chunks = self._create_chunks(content, filename)
        print(chunks)
        return chunks
    
    async def _process_pdf(self, file_path: str) -> str:
        """Extract text from PDF including OCR for images"""
        text = ""
        doc = fitz.open(file_path)
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            
            # Extract text
            page_text = page.get_text()
            
            # If no text found, try OCR on images
            if not page_text.strip():
                pix = page.get_pixmap()
                img_data = pix.tobytes("png")
                image = Image.open(io.BytesIO(img_data))
                page_text = pytesseract.image_to_string(image)
            
            text += f"\n--- Page {page_num + 1} ---\n{page_text}"
        
        doc.close()
        return text
    
    async def _process_docx(self, file_path: str) -> str:
        """Extract text from DOCX"""
        doc = docx.Document(file_path)
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        return text
    
    async def _process_txt(self, file_path: str) -> str:
        """Process text file"""
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()
    
    async def _process_image(self, file_path: str) -> str:
        """Extract text from image using OCR"""
        image = Image.open(file_path)
        text = pytesseract.image_to_string(image)
        return text
    
    def _create_chunks(self, content: str, filename: str) -> List[DocumentChunk]:
        """Split content into chunks"""
        from app.utils.config import settings
        
        chunks = []
        words = content.split()
        
        for i in range(0, len(words), settings.CHUNK_SIZE - settings.CHUNK_OVERLAP):
            chunk_words = words[i:i + settings.CHUNK_SIZE]
            chunk_content = " ".join(chunk_words)
            
            chunk_id = hashlib.md5(f"{filename}_{i}".encode()).hexdigest()
            
            chunk = DocumentChunk(
                chunk_id=chunk_id,
                content=chunk_content,
                metadata={
                    "filename": filename,
                    "chunk_index": i // (settings.CHUNK_SIZE - settings.CHUNK_OVERLAP),
                    "word_count": len(chunk_words)
                }
            )
            chunks.append(chunk)
        
        return chunks
